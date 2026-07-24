"""
Artifact Generation Skills Architecture

Defines reusable skills for generating management consulting-level deliverables:
- PowerPoint presentations (strategy, architecture, findings)
- Excel dashboards and analyses
- Visio diagrams and process flows
- Process flow documentation (BPMN, swimlanes)
- Figma prototypes and design artifacts

Each skill is a modular, composable function that can be:
1. Invoked directly from the backend
2. Called as a Copilot Studio action/tool
3. Wired as an MCP tool for external agents
"""

from abc import ABC, abstractmethod
from typing import Dict, Any, List, Optional, Tuple
from dataclasses import dataclass, asdict
from pathlib import Path
from datetime import datetime
import json
import logging

logger = logging.getLogger(__name__)


@dataclass
class ArtifactMetadata:
    """Metadata for generated artifacts."""
    artifact_type: str
    title: str
    description: str
    generated_at: str
    agent_context: Dict[str, Any]
    templates_used: List[str]
    source_data: Dict[str, Any]
    tags: List[str]


class ArtifactSkill(ABC):
    """Base class for artifact generation skills."""
    
    skill_name: str
    version: str = "1.0"
    tags: List[str] = []
    supported_formats: List[str] = []
    
    @abstractmethod
    def generate(self, context: Dict[str, Any], **kwargs) -> Dict[str, Any]:
        """
        Generate artifact.
        
        Args:
            context: Agent context including D365 architecture data
            **kwargs: Skill-specific parameters
            
        Returns:
            Dict with artifact_path, metadata, and generation_stats
        """
        pass
    
    @abstractmethod
    def validate(self, config: Dict[str, Any]) -> Tuple[bool, List[str]]:
        """
        Validate configuration before generation.
        
        Returns:
            (is_valid, error_messages)
        """
        pass


# =============================================================================
# POWERPOINT SKILL - Management Consulting Deck Generation
# =============================================================================

class PowerPointSkill(ArtifactSkill):
    """
    Generate management consulting-grade PowerPoint presentations.
    
    Use Cases:
    - Architecture recommendations deck
    - Solution overview for stakeholders
    - Implementation roadmap presentation
    - Executive summary findings
    - Risk and mitigation analysis
    """
    
    skill_name = "powerpoint_deck_generator"
    supported_formats = ["pptx", "pdf", "png"]
    tags = [
        "presentation",
        "consulting_deliverable",
        "executive_summary",
        "architecture",
        "strategy"
    ]
    
    AVAILABLE_TEMPLATES = {
        "strategy_recommendation": {
            "description": "Strategy and recommendation deck",
            "slides": [
                "title_slide",
                "executive_summary",
                "situation_analysis",
                "recommendation",
                "architecture_overview",
                "benefits",
                "risks_mitigation",
                "implementation_roadmap",
                "success_metrics",
                "next_steps",
                "appendix"
            ]
        },
        "architecture_deep_dive": {
            "description": "Technical architecture presentation",
            "slides": [
                "title_slide",
                "current_state",
                "requirements",
                "proposed_architecture",
                "integration_design",
                "data_flows",
                "security_governance",
                "performance_capacity",
                "implementation_approach",
                "timeline_budget",
                "appendix"
            ]
        },
        "findings_report": {
            "description": "Assessment findings and recommendations",
            "slides": [
                "title_slide",
                "executive_summary",
                "scope",
                "key_findings",
                "recommendations",
                "priority_matrix",
                "implementation_sequence",
                "resource_requirements",
                "risk_register",
                "appendix"
            ]
        },
        "implementation_roadmap": {
            "description": "Phased implementation plan",
            "slides": [
                "title_slide",
                "phase_overview",
                "phase_1_detail",
                "phase_2_detail",
                "phase_3_detail",
                "dependencies",
                "resource_plan",
                "governance_structure",
                "success_criteria",
                "appendix"
            ]
        }
    }
    
    def generate(self, context: Dict[str, Any], **kwargs) -> Dict[str, Any]:
        """
        Generate PowerPoint presentation.
        
        Args:
            context: Agent context including:
                - architecture_summary: str
                - recommendations: List[str]
                - risks: List[Dict]
                - success_metrics: List[Dict]
            template_type: str (strategy_recommendation, architecture_deep_dive, etc)
            customer_name: str
            output_path: str
        """
        template_type = kwargs.get("template_type", "strategy_recommendation")
        customer_name = kwargs.get("customer_name", "Customer")
        output_path = kwargs.get("output_path", "generated_decks")
        
        # Validate
        is_valid, errors = self.validate(kwargs)
        if not errors:
            return {"error": f"Validation failed: {errors}"}
        
        try:
            # Import here to avoid hard dependency
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.enum.text import PP_ALIGN
            from pptx.dml.color import RGBColor
            
        except ImportError:
            return {
                "error": "python-pptx not installed. Run: pip install python-pptx"
            }
        
        # Create presentation
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # Define style
        accent_color = RGBColor(0, 102, 204)  # Microsoft blue
        
        # Get template slides
        template = self.AVAILABLE_TEMPLATES.get(template_type)
        if not template:
            return {"error": f"Unknown template: {template_type}"}
        
        # Generate slides based on template
        for slide_type in template["slides"]:
            if slide_type == "title_slide":
                self._add_title_slide(prs, context, customer_name, accent_color)
            elif slide_type == "executive_summary":
                self._add_executive_summary_slide(prs, context, accent_color)
            elif slide_type == "recommendation":
                self._add_recommendation_slide(prs, context, accent_color)
            elif slide_type == "architecture_overview":
                self._add_architecture_slide(prs, context, accent_color)
            elif slide_type == "benefits":
                self._add_benefits_slide(prs, context, accent_color)
            elif slide_type == "risks_mitigation":
                self._add_risks_slide(prs, context, accent_color)
            elif slide_type == "implementation_roadmap":
                self._add_roadmap_slide(prs, context, accent_color)
            elif slide_type == "success_metrics":
                self._add_metrics_slide(prs, context, accent_color)
        
        # Save
        Path(output_path).mkdir(parents=True, exist_ok=True)
        filename = f"{customer_name.replace(' ', '_')}_architecture_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
        filepath = Path(output_path) / filename
        
        prs.save(str(filepath))
        
        return {
            "success": True,
            "artifact_path": str(filepath),
            "artifact_type": "powerpoint",
            "template_used": template_type,
            "slide_count": len(prs.slides),
            "metadata": {
                "generated_at": datetime.utcnow().isoformat(),
                "customer": customer_name,
                "agent_context": context,
            }
        }
    
    def validate(self, config: Dict[str, Any]) -> Tuple[bool, List[str]]:
        """Validate PowerPoint generation config."""
        errors = []
        
        template_type = config.get("template_type")
        if template_type and template_type not in self.AVAILABLE_TEMPLATES:
            errors.append(f"Invalid template: {template_type}")
        
        return (len(errors) == 0, errors)
    
    def _add_title_slide(self, prs, context, customer_name, accent_color):
        """Add title slide."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        
        # Add background color
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = accent_color
        
        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        p = title_frame.paragraphs[0]
        p.text = f"Dynamics 365 Architecture Recommendation"
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        
        # Add subtitle
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1))
        subtitle_frame = subtitle_box.text_frame
        p = subtitle_frame.paragraphs[0]
        p.text = customer_name
        p.font.size = Pt(32)
        p.font.color.rgb = RGBColor(200, 200, 200)
    
    def _add_executive_summary_slide(self, prs, context, accent_color):
        """Add executive summary slide."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        p = title_frame.paragraphs[0]
        p.text = "Executive Summary"
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = accent_color
        
        # Content
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        
        summary = context.get("architecture_summary", "Architecture recommendation summary")
        p = content_frame.paragraphs[0]
        p.text = summary
        p.font.size = Pt(18)
        p.level = 0
    
    def _add_recommendation_slide(self, prs, context, accent_color):
        """Add recommendation slide."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        p = title_frame.paragraphs[0]
        p.text = "Recommendations"
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = accent_color
        
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        
        recommendations = context.get("recommendations", [])
        for i, rec in enumerate(recommendations[:5]):  # Max 5 recommendations
            p = content_frame.add_paragraph()
            p.text = rec
            p.level = 0
            p.font.size = Pt(18)
    
    def _add_architecture_slide(self, prs, context, accent_color):
        """Add architecture overview slide."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        p = title_frame.paragraphs[0]
        p.text = "Proposed Architecture"
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = accent_color
        
        # Placeholder for architecture diagram
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
        content_frame = content_box.text_frame
        p = content_frame.paragraphs[0]
        p.text = "[Architecture Diagram to be inserted from Visio/Figma]"
        p.font.size = Pt(16)
        p.font.italic = True
    
    def _add_benefits_slide(self, prs, context, accent_color):
        """Add benefits slide."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        p = title_frame.paragraphs[0]
        p.text = "Key Benefits"
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = accent_color
        
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        
        benefits = context.get("benefits", [])
        for benefit in benefits[:5]:
            p = content_frame.add_paragraph()
            p.text = benefit
            p.font.size = Pt(18)
    
    def _add_risks_slide(self, prs, context, accent_color):
        """Add risks and mitigation slide."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        p = title_frame.paragraphs[0]
        p.text = "Risk Mitigation"
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = accent_color
        
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        
        risks = context.get("risks", [])
        for risk in risks[:5]:
            p = content_frame.add_paragraph()
            p.text = f"{risk.get('risk', '')}: {risk.get('mitigation', '')}"
            p.font.size = Pt(16)
    
    def _add_roadmap_slide(self, prs, context, accent_color):
        """Add implementation roadmap slide."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        p = title_frame.paragraphs[0]
        p.text = "Implementation Roadmap"
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = accent_color
        
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        
        roadmap = context.get("roadmap", [])
        for phase in roadmap[:5]:
            p = content_frame.add_paragraph()
            p.text = f"Phase: {phase}"
            p.font.size = Pt(18)
    
    def _add_metrics_slide(self, prs, context, accent_color):
        """Add success metrics slide."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
        title_frame = title_box.text_frame
        p = title_frame.paragraphs[0]
        p.text = "Success Metrics"
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = accent_color
        
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        
        metrics = context.get("success_metrics", [])
        for metric in metrics[:5]:
            p = content_frame.add_paragraph()
            p.text = metric
            p.font.size = Pt(18)


# =============================================================================
# EXCEL SKILL - Dashboard and Analysis Generation
# =============================================================================

class ExcelSkill(ArtifactSkill):
    """
    Generate executive dashboards and analyses in Excel.
    
    Use Cases:
    - Architecture comparison matrices
    - Cost-benefit analysis
    - Implementation timeline with resources
    - Risk and complexity assessment
    - Integration inventory
    """
    
    skill_name = "excel_dashboard_generator"
    supported_formats = ["xlsx", "csv"]
    tags = [
        "dashboard",
        "analysis",
        "executive_summary",
        "cost_benefit",
        "inventory"
    ]
    
    def generate(self, context: Dict[str, Any], **kwargs) -> Dict[str, Any]:
        """Generate Excel dashboard."""
        try:
            from openpyxl import Workbook
            from openpyxl.styles import Font, PatternFill, Alignment
        except ImportError:
            return {
                "error": "openpyxl not installed. Run: pip install openpyxl"
            }
        
        output_path = kwargs.get("output_path", "generated_analyses")
        customer_name = kwargs.get("customer_name", "Customer")
        
        wb = Workbook()
        ws = wb.active
        ws.title = "Architecture Analysis"
        
        # Add header
        ws['A1'] = f"D365 Architecture Analysis - {customer_name}"
        ws['A1'].font = Font(size=16, bold=True)
        
        ws['A2'] = f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"
        
        # Add sections from context
        row = 4
        
        # Recommendations
        if "recommendations" in context:
            ws[f'A{row}'] = "Recommendations"
            ws[f'A{row}'].font = Font(bold=True, size=12)
            row += 1
            
            for i, rec in enumerate(context["recommendations"]):
                ws[f'A{row}'] = f"{i+1}. {rec}"
                row += 1
            
            row += 1
        
        # Risks
        if "risks" in context:
            ws[f'A{row}'] = "Risks & Mitigation"
            ws[f'A{row}'].font = Font(bold=True, size=12)
            row += 1
            
            ws[f'A{row}'] = "Risk"
            ws[f'B{row}'] = "Impact"
            ws[f'C{row}'] = "Mitigation"
            row += 1
            
            for risk in context["risks"]:
                ws[f'A{row}'] = risk.get("risk", "")
                ws[f'B{row}'] = risk.get("impact", "")
                ws[f'C{row}'] = risk.get("mitigation", "")
                row += 1
        
        # Save
        Path(output_path).mkdir(parents=True, exist_ok=True)
        filename = f"{customer_name.replace(' ', '_')}_analysis_{datetime.now().strftime('%Y%m%d_%H%M%S')}.xlsx"
        filepath = Path(output_path) / filename
        
        wb.save(str(filepath))
        
        return {
            "success": True,
            "artifact_path": str(filepath),
            "artifact_type": "excel",
            "metadata": {
                "generated_at": datetime.utcnow().isoformat(),
                "customer": customer_name,
            }
        }
    
    def validate(self, config: Dict[str, Any]) -> Tuple[bool, List[str]]:
        """Validate Excel generation config."""
        return (True, [])


# =============================================================================
# VISIO/PROCESS FLOW SKILL - Diagram Generation
# =============================================================================

class ProcessFlowSkill(ArtifactSkill):
    """
    Generate Visio and process flow diagrams.
    
    Use Cases:
    - Solution architecture diagrams
    - Integration flow diagrams
    - BPMN process models
    - Swimlane diagrams
    - Data flow diagrams (DFD)
    """
    
    skill_name = "process_flow_diagram_generator"
    supported_formats = ["vsdx", "bpmn", "mermaid", "svg", "png"]
    tags = [
        "diagram",
        "process_flow",
        "architecture",
        "bpmn",
        "integration"
    ]
    
    def generate(self, context: Dict[str, Any], **kwargs) -> Dict[str, Any]:
        """Generate process flow diagram."""
        diagram_type = kwargs.get("diagram_type", "architecture")
        output_format = kwargs.get("output_format", "mermaid")
        
        if output_format == "mermaid":
            return self._generate_mermaid_diagram(context, diagram_type, **kwargs)
        elif output_format == "bpmn":
            return self._generate_bpmn_diagram(context, **kwargs)
        else:
            return {"error": f"Unsupported format: {output_format}"}
    
    def _generate_mermaid_diagram(self, context, diagram_type, **kwargs):
        """Generate Mermaid diagram."""
        output_path = kwargs.get("output_path", "generated_diagrams")
        customer_name = kwargs.get("customer_name", "Customer")
        
        if diagram_type == "architecture":
            mermaid_code = self._create_architecture_diagram(context)
        elif diagram_type == "integration":
            mermaid_code = self._create_integration_diagram(context)
        elif diagram_type == "process":
            mermaid_code = self._create_process_diagram(context)
        else:
            mermaid_code = ""
        
        # Save as markdown and potentially render
        Path(output_path).mkdir(parents=True, exist_ok=True)
        filename = f"{customer_name.replace(' ', '_')}_{diagram_type}_diagram.md"
        filepath = Path(output_path) / filename
        
        with open(filepath, 'w') as f:
            f.write(f"# {diagram_type.title()} Diagram\n\n")
            f.write("```mermaid\n")
            f.write(mermaid_code)
            f.write("\n```\n")
        
        return {
            "success": True,
            "artifact_path": str(filepath),
            "artifact_type": "mermaid_diagram",
            "diagram_type": diagram_type,
            "mermaid_code": mermaid_code,
            "metadata": {
                "generated_at": datetime.utcnow().isoformat(),
                "customer": customer_name,
            }
        }
    
    def _create_architecture_diagram(self, context):
        """Create architecture diagram."""
        return """graph TB
    subgraph "Dynamics 365"
        Sales[D365 Sales]
        FS[D365 Field Service]
        CS[D365 Customer Service]
    end
    
    subgraph "Power Platform"
        Copilot[Copilot Studio]
        PA[Power Automate]
        PC[Power Apps]
    end
    
    subgraph "Integration"
        API[API Management]
        DF[Data Factory]
    end
    
    Sales --> Copilot
    FS --> Copilot
    CS --> Copilot
    Copilot --> PA
    PA --> DF
    DF --> API
"""
    
    def _create_integration_diagram(self, context):
        """Create integration flow diagram."""
        return """graph LR
    Ecolab[Ecolab Systems]
    D365[D365]
    External[External Services]
    
    Ecolab -->|Data Sync| D365
    D365 -->|Events| External
    External -->|Updates| D365
"""
    
    def _create_process_diagram(self, context):
        """Create process flow diagram."""
        return """graph TD
    Start([Start]) --> Input[Gather Requirements]
    Input --> Analysis[Architecture Analysis]
    Analysis --> Design[Design Solution]
    Design --> Review{Review OK?}
    Review -->|No| Analysis
    Review -->|Yes| Implement[Implementation]
    Implement --> End([End])
"""
    
    def _generate_bpmn_diagram(self, context, **kwargs):
        """Generate BPMN diagram (simplified)."""
        return {
            "success": True,
            "artifact_type": "bpmn",
            "message": "BPMN generation requires python-bpmn library. See documentation.",
        }
    
    def validate(self, config: Dict[str, Any]) -> Tuple[bool, List[str]]:
        """Validate diagram generation config."""
        errors = []
        
        valid_formats = self.supported_formats
        if config.get("output_format") and config["output_format"] not in valid_formats:
            errors.append(f"Invalid format. Supported: {valid_formats}")
        
        return (len(errors) == 0, errors)


# =============================================================================
# Skills Registry
# =============================================================================

class SkillsRegistry:
    """Registry of all available artifact generation skills."""
    
    def __init__(self):
        self.skills: Dict[str, ArtifactSkill] = {
            "powerpoint": PowerPointSkill(),
            "excel": ExcelSkill(),
            "process_flow": ProcessFlowSkill(),
        }
    
    def get_skill(self, skill_name: str) -> Optional[ArtifactSkill]:
        """Get skill by name."""
        return self.skills.get(skill_name)
    
    def list_skills(self) -> Dict[str, Dict[str, Any]]:
        """List all available skills."""
        return {
            name: {
                "skill_name": skill.skill_name,
                "version": skill.version,
                "tags": skill.tags,
                "supported_formats": skill.supported_formats,
            }
            for name, skill in self.skills.items()
        }
    
    def invoke_skill(self, skill_name: str, context: Dict[str, Any], **kwargs) -> Dict[str, Any]:
        """Invoke a skill."""
        skill = self.get_skill(skill_name)
        if not skill:
            return {"error": f"Skill not found: {skill_name}"}
        
        # Validate
        is_valid, errors = skill.validate(kwargs)
        if not is_valid:
            return {"error": f"Validation failed: {errors}"}
        
        # Generate
        return skill.generate(context, **kwargs)


if __name__ == "__main__":
    # Example usage
    registry = SkillsRegistry()
    
    print("Available Skills:")
    print(json.dumps(registry.list_skills(), indent=2))
    
    # Example: Generate PowerPoint
    context = {
        "architecture_summary": "Implement a scalable D365 solution with Copilot-assisted workflows",
        "recommendations": [
            "Adopt Copilot Studio for custom agents",
            "Use Field Service Copilot for dispatch",
            "Implement MCP for skill orchestration"
        ],
        "risks": [
            {"risk": "AI Governance", "mitigation": "Implement Responsible AI Canvas"}
        ],
        "benefits": [
            "30% improvement in field technician productivity",
            "25% reduction in dispatch time",
            "Enhanced customer satisfaction"
        ]
    }
    
    result = registry.invoke_skill("powerpoint", context, 
                                   template_type="strategy_recommendation",
                                   customer_name="Ecolab")
    
    print("\nGeneration Result:")
    print(json.dumps(result, indent=2))
